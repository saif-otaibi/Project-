from pptx import Presentation
from docx import Document
from fpdf import FPDF
from win32_setctime import setctime
from pathlib import Path
import os
import subprocess
import time
import hashlib

hasher = hashlib.sha256()
pattern = '%d.%m.%Y'

def run(cmd):
    completed = subprocess.Popen("powershell -command " +cmd , stdin=subprocess.PIPE,stdout=subprocess.PIPE, shell=True)
    stdout_value = completed.communicate()[0]
    return stdout_value

def getdate():
    global year,month,day
    year = 0
    while (year > 2040) or (year < 1980):
        try:
            year = int(input("Enter Correct Year Between 1980 - 2040: "))
        except:
            print("This field accepts numbers only.")
    month = 0
    while (month > 12) or (month < 1):
        try:
            month = int(input("Enter Correct Month Between 1 - 12: "))
        except:
            print("This field accepts numbers only.")
    if month in (1,3,5,7,8,10,12):
        day = 0
        while (day > 31) or (day < 1):
            try:
                day = int(input("Enter Correct Day Between 1 - 31: "))
            except:
                print("This field accepts numbers only.")
    elif month in (4,6,9,11):
        day = 0
        while (day > 30) or (day < 1):
            try:
                day = int(input("Enter Correct Day Between 1 - 30: "))
            except:
                print("This field accepts numbers only.")
    elif month == 2 and year % 4 == 0:
        day = 0
        while (day > 29) or (day < 1):
            try:
                day = int(input("Enter Correct Day Between 1 - 29: "))
            except:
                print("This field accepts numbers only.")
    elif month == 2 and year % 4 != 0:
        day = 0
        while (day > 28) or (day < 1):
            try:
                day = int(input("Enter Correct Day Between 1 - 28: "))
            except:
                print("This field accepts numbers only.")

print("Welcome to Capastone2")
print()
folder = ""
while (folder == ""):
    tmp = input("Enter location to work in (or empty for current location): ").replace('\\','/')
    if tmp == "":
        break
    p = Path(tmp)
    if p.is_dir():
        folder = tmp
    else:
        print(tmp + " is a wrong location" )
if folder != "":
    os.chdir(folder)


print()

print("1) PPTX")
print("2) DOCX")
print("3) TXT")
print("4) PDF")
filetype = ""
while not (filetype == "1" or filetype == "2" or filetype == "3" or filetype == "4"):
    filetype = input("Please select 1,2,3 or 4: ")


if(filetype == "1"):
    print("PPTX selected.")
    name = ""
    while (name == ""):
        name = input("Enter filename: ")
    prs=Presentation()
    lyt=prs.slide_layouts[0]
    slide=prs.slides.add_slide(lyt)
    title=slide.shapes.title
    subtitle=slide.placeholders[1]
    title.text="Hey,This is a Slide! How exciting!"
    subtitle.text="Really?"
    prs.save(name + ".pptx")
    ifmod = ""
    while not ((ifmod == "Y") or (ifmod == "y") or (ifmod == "N") or (ifmod == "n")):
        ifmod=input("Do you want to enter cereation date? [Y/N] ")
    if ifmod == "Y" or ifmod == "y":
        print("Please Enter cereation date: ")
        getdate()
        cyear=year
        cmonth=month
        cday=day
        cdate = str(cday) + '.' + str(cmonth) + '.' + str(cyear)
        cepoch = int(time.mktime(time.strptime(cdate, pattern)))
        setctime(name + ".pptx", cepoch)
        
    ifmod = ""
    while not ((ifmod == "Y") or (ifmod == "y") or (ifmod == "N") or (ifmod == "n")):
        ifmod=input("Do you want to enter modification date? [Y/N] ")
    if ifmod == "Y" or ifmod == "y":
        print("Please Enter modification date: ")
        getdate()
        myear=year
        mmonth=month
        mday=day
        mdate = str(mday) + '.' + str(mmonth) + '.' + str(myear)
        mepoch = int(time.mktime(time.strptime(mdate, pattern)))
        os.utime(name + ".pptx", (mepoch, mepoch))
    with open(name + '.pptx','rb') as open_file:
        content = open_file.read()
        hasher.update(content)
    print ('The hash of our file is:\n'+ hasher.hexdigest())
    ifmod = ""
    while not ((ifmod == "Y") or (ifmod == "y") or (ifmod == "N") or (ifmod == "n")):
        ifmod=input("Do you want to enter access date? [Y/N] ")
    if ifmod == "Y" or ifmod == "y":
        print("Please Enter access date: ")
        getdate()
        ayear=year
        amonth=month
        aday=day
        apath= os.getcwd() + '\\' + name + ".pptx"
        powershell = "(Get-Item '" + str(apath) + "').LastAccessTime=('" + str(amonth) + " " + str(aday) + " " + str(ayear) + " 00:00:00')"
        print(powershell)
        run(powershell)


if(filetype == "2"):
    print("DOCX selected.")
    name = ""
    while (name == ""):
        name = input("Enter filename: ")
    document = Document()
    document.add_heading('This is a file for all my passwords and usernames', 0)
    p = document.add_paragraph(
             "Microsoft Office Username=sageab123, Pass=123456abcd\n"
             "Youtube  Username=sageab_goo, Pass=go123@sageab\n")
    p.add_run(' some bold text').bold = True
    p.add_run(' and italic text.').italic = True
    p = document.add_paragraph('And this is another text paragraph')
    document.save(name + '.docx')
    ifmod = ""
    while not ((ifmod == "Y") or (ifmod == "y") or (ifmod == "N") or (ifmod == "n")):
        ifmod=input("Do you want to enter cereation date? [Y/N] ")
    if ifmod == "Y" or ifmod == "y":
        print("Please Enter cereation date: ")
        getdate()
        cyear=year
        cmonth=month
        cday=day
        cdate = str(cday) + '.' + str(cmonth) + '.' + str(cyear)
        cepoch = int(time.mktime(time.strptime(cdate, pattern)))
        setctime(name + ".docx", cepoch)
        
    ifmod = ""
    while not ((ifmod == "Y") or (ifmod == "y") or (ifmod == "N") or (ifmod == "n")):
        ifmod=input("Do you want to enter modification date? [Y/N] ")
    if ifmod == "Y" or ifmod == "y":
        print("Please Enter modification date: ")
        getdate()
        myear=year
        mmonth=month
        mday=day
        mdate = str(mday) + '.' + str(mmonth) + '.' + str(myear)
        mepoch = int(time.mktime(time.strptime(mdate, pattern)))
        os.utime(name + ".docx", (mepoch, mepoch))
    with open(name + '.docx','rb') as open_file:
        content = open_file.read()
        hasher.update(content)
    print ('The hash of our file is:\n'+ hasher.hexdigest())
    ifmod = ""
    while not ((ifmod == "Y") or (ifmod == "y") or (ifmod == "N") or (ifmod == "n")):
        ifmod=input("Do you want to enter access date? [Y/N] ")
    if ifmod == "Y" or ifmod == "y":
        print("Please Enter access date: ")
        getdate()
        ayear=year
        amonth=month
        aday=day
        apath= os.getcwd() + '\\' + name + ".docx"
        powershell = "(Get-Item '" + str(apath) + "').LastAccessTime=('" + str(amonth) + " " + str(aday) + " " + str(ayear) + " 00:00:00')"
        print(powershell)
        run(powershell)

if(filetype == "3"):
    print("TXT selected.")
    name = ""
    while (name == ""):
        name = input("Enter filename: ")
    fh = open(name + '.txt', 'w')
    fh.write("This is a file for all my passwords and usernames\n"
             "Microsoft Office Username=sageab123, Pass=123456abcd\n"
             "Youtube  Username=sageab_goo, Pass=go123@sageab\n")
    fh.close()
    ifmod = ""
    while not ((ifmod == "Y") or (ifmod == "y") or (ifmod == "N") or (ifmod == "n")):
        ifmod=input("Do you want to enter cereation date? [Y/N] ")
    if ifmod == "Y" or ifmod == "y":
        print("Please Enter cereation date: ")
        getdate()
        cyear=year
        cmonth=month
        cday=day
        cdate = str(cday) + '.' + str(cmonth) + '.' + str(cyear)
        cepoch = int(time.mktime(time.strptime(cdate, pattern)))
        setctime(name + ".txt", cepoch)
        
    ifmod = ""
    while not ((ifmod == "Y") or (ifmod == "y") or (ifmod == "N") or (ifmod == "n")):
        ifmod=input("Do you want to enter modification date? [Y/N] ")
    if ifmod == "Y" or ifmod == "y":
        print("Please Enter modification date: ")
        getdate()
        myear=year
        mmonth=month
        mday=day
        mdate = str(mday) + '.' + str(mmonth) + '.' + str(myear)
        mepoch = int(time.mktime(time.strptime(mdate, pattern)))
        os.utime(name + ".txt", (mepoch, mepoch))
    with open(name + '.txt','rb') as open_file:
        content = open_file.read()
        hasher.update(content)
    print ('The hash of our file is:\n'+ hasher.hexdigest())
    ifmod = ""
    while not ((ifmod == "Y") or (ifmod == "y") or (ifmod == "N") or (ifmod == "n")):
        ifmod=input("Do you want to enter access date? [Y/N] ")
    if ifmod == "Y" or ifmod == "y":
        print("Please Enter access date: ")
        getdate()
        ayear=year
        amonth=month
        aday=day
        apath= os.getcwd() + '\\' + name + ".txt"
        powershell = "(Get-Item '" + str(apath) + "').LastAccessTime=('" + str(amonth) + " " + str(aday) + " " + str(ayear) + " 00:00:00')"
        print(powershell)
        run(powershell)

if(filetype == "4"):
    print("PDF selected.")
    name = ""
    while (name == ""):
        name = input("Enter filename: ")
    pdf = FPDF('P', 'mm', 'A4') 
    pdf.add_page()
    pdf.set_left_margin(2.4)
    pdf.image('C:\\Users\\DELLL\OneDrive\\Desktop\\schedule2.PNG',25.4,60,160) 
    pdf.set_font('Times', '', 16)
    pdf.set_xy(25.4, 140)  

    #cell creates a single line text box
    pdf.cell(170, 30, 'My bicycle fell because it was two-tired.', 1)
    pdf.set_xy(25.4, 180)
    pdf.set_text_color(255, 0, 0)

    #multi_cell creates multiple lines text box
    pdf.multi_cell(170, 10, 'Hello World!! this is just non-sense to see whether the text will go to the next line or continue out of page as in a singular cell. I think this works perfectly.', 0, 'L', 0)
    pdf.set_author("Abdulrahim")
    pdf.set_creator('Abdu')
    pdf.output(name + '.pdf', 'F')
    ifmod = ""
    while not ((ifmod == "Y") or (ifmod == "y") or (ifmod == "N") or (ifmod == "n")):
        ifmod=input("Do you want to enter cereation date? [Y/N] ")
    if ifmod == "Y" or ifmod == "y":
        print("Please Enter cereation date: ")
        getdate()
        cyear=year
        cmonth=month
        cday=day
        cdate = str(cday) + '.' + str(cmonth) + '.' + str(cyear)
        cepoch = int(time.mktime(time.strptime(cdate, pattern)))
        setctime(name + ".pdf", cepoch)
        
    ifmod = ""
    while not ((ifmod == "Y") or (ifmod == "y") or (ifmod == "N") or (ifmod == "n")):
        ifmod=input("Do you want to enter modification date? [Y/N] ")
    if ifmod == "Y" or ifmod == "y":
        print("Please Enter modification date: ")
        getdate()
        myear=year
        mmonth=month
        mday=day
        mdate = str(mday) + '.' + str(mmonth) + '.' + str(myear)
        mepoch = int(time.mktime(time.strptime(mdate, pattern)))
        os.utime(name + ".pdf", (mepoch, mepoch))
    with open(name + '.pdf','rb') as open_file:
        content = open_file.read()
        hasher.update(content)
    print ('The hash of our file is:\n'+ hasher.hexdigest())
    ifmod = ""
    while not ((ifmod == "Y") or (ifmod == "y") or (ifmod == "N") or (ifmod == "n")):
        ifmod=input("Do you want to enter access date? [Y/N] ")
    if ifmod == "Y" or ifmod == "y":
        print("Please Enter access date: ")
        getdate()
        ayear=year
        amonth=month
        aday=day
        apath= os.getcwd() + '\\' + name + ".pdf"
        powershell = "(Get-Item '" + str(apath) + "').LastAccessTime=('" + str(amonth) + " " + str(aday) + " " + str(ayear) + " 00:00:00')"
        print(powershell)
        run(powershell)

    
